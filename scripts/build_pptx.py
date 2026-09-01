#!/usr/bin/env python3
"""Place finished slide images into a wide PowerPoint without overlays."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("input", type=Path, help="Directory containing final PNG/JPEG slides")
    parser.add_argument("output", type=Path, help="Output PPTX path")
    args = parser.parse_args()

    slides = sorted([*args.input.glob("*.png"), *args.input.glob("*.jpg"), *args.input.glob("*.jpeg")])
    if not slides:
        raise SystemExit(f"No slide images found in {args.input}")

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    for image_path in slides:
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

    args.output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(args.output)
    print(f"built {args.output} from {len(slides)} slides")


if __name__ == "__main__":
    main()
